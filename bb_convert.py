import cv2
from lxml import etree
import os
from PIL import Image
from pptx import Presentation
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.util import Inches

from x_utils import *

THUMBNAIL_FILENAME = f"{PATH}/thumbnail.png"
VIDEO_FILENAME = f"{PATH}/video.mov"

def convert():
    height = width = None

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    pages = []

    video_nr = 0
    for frames in read_output_videos(verbose=True):
        if height is None:
            height, width, _ = frames[0].shape

        pages.append(Image.fromarray(cv2.cvtColor(frames[-1], cv2.COLOR_BGR2RGB)))

        cv2.imwrite(THUMBNAIL_FILENAME, frames[0])
        video_writer = cv2.VideoWriter(VIDEO_FILENAME, cv2.VideoWriter_fourcc(*"mp4v"), FINAL_FRAMERATE, (width, height))
        for frame in frames:
            video_writer.write(frame)
        video_writer.release()

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        movie = slide.shapes.add_movie(VIDEO_FILENAME, 0, 0, prs.slide_width, prs.slide_height, poster_frame_image=THUMBNAIL_FILENAME)
        tree = movie._element.getparent().getparent().getnext().getnext()
        timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == "cond"][0]
        timing.set("delay", "0")

        print(f"\033[30;1mProcessed video #{video_nr + 1}\033[0m")

        video_nr += 1

    pages[0].save(f"{PATH}/output.pdf", "PDF", resolution=100.0, save_all=True, append_images=pages[1:])

    prs.save(f"{PATH}/output.pptx")
    os.unlink(VIDEO_FILENAME)
    os.unlink(THUMBNAIL_FILENAME)

    print(f"\033[32;1mDone!\033[0m")

if __name__ == "__main__":
    convert()
